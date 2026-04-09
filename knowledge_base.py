import os
import re
import io
import json
import tempfile
from pathlib import Path
from typing import Dict, List, Optional, Callable
from datetime import datetime

import numpy as np
import streamlit as st
import pandas as pd
import oss2
from oss2.exceptions import NoSuchKey

from config import AppConfig
from oss_client import get_oss_bucket
from project_manager import ProjectManager
from vector_store import SimpleVectorStore

# 可选依赖标志和导入
HAS_DOCX = False
try:
    import docx
    HAS_DOCX = True
except ImportError:
    pass

HAS_PDF = False
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    pass

HAS_EXCEL = False
try:
    from openpyxl import load_workbook
    HAS_EXCEL = True
except ImportError:
    load_workbook = None

HAS_PPTX = False
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    Presentation = None

HAS_XMIND = False
try:
    import xmindparser
    HAS_XMIND = True
except ImportError:
    xmindparser = None

# ====================== 知识库封装（OSS 存储） ======================

class EnhancedKnowledgeBase:
    def __init__(self, project_name: str):
        self.project_name = project_name
        self.oss_bucket = get_oss_bucket()
        self.kb_prefix = ProjectManager.get_kb_path(project_name)  # OSS 前缀
        self.vector_prefix = ProjectManager.get_vector_store_path(project_name)
        # 临时目录
        self.temp_dir = Path(tempfile.gettempdir()) / f"kb_{project_name}"
        self.temp_dir.mkdir(parents=True, exist_ok=True)

        self.vector_store = SimpleVectorStore(
            str(self.temp_dir / "vector_store"),
            self.oss_bucket,
            self.vector_prefix
        )
        self.chunk_size = 800
        self.chunk_overlap = 200
        self.index_loaded = False
        self._auto_load()
    def _auto_load(self):
        files = self.get_file_list(with_metadata=True)
        if files and self.vector_store.is_index_valid(files):
            if self.vector_store.load_index():
                self.index_loaded = True

    def refresh_index(self) -> Dict:
        self.index_loaded = False
        files = self.get_file_list(with_metadata=True)
        if not files:
            self.vector_store.clear_index()
            return {"status": "error", "message": "知识库无文件"}
        if self.vector_store.load_index():
            self.index_loaded = True
            stats = self.vector_store.get_stats()
            built_files = self.vector_store.get_built_files()
            all_files = {f["name"] for f in files}
            pending = all_files - built_files
            if pending:
                msg = f"已加载索引，包含 {stats['documents']} 个片段。注意：有 {len(pending)} 个文件未构建，请点击「构建知识库」增量添加。"
            else:
                msg = f"已加载索引，包含 {stats['documents']} 个知识片段"
            return {"status": "success", "message": msg}
        return {"status": "error", "message": "未找到有效索引，请先构建知识库"}

    def get_file_list(self, with_metadata=False):
        files = []
        prefix = f"{self.kb_prefix}/"
        try:
            for obj in oss2.ObjectIterator(self.oss_bucket, prefix=prefix):
                key = obj.key
                filename = key[len(prefix):]
                if not filename or filename.startswith('.') or filename.endswith('/'):
                    continue
                ext = filename.split('.')[-1].lower()
                if ext in AppConfig.ALLOWED_FILE_TYPES:
                    size_kb = obj.size / 1024
                    info = {"name": filename, "size": f"{size_kb:.1f}KB"}
                    if with_metadata:
                        # 兼容 last_modified 可能是 datetime 或 时间戳
                        if hasattr(obj.last_modified, 'timestamp'):
                            mtime = obj.last_modified.timestamp()
                        else:
                            mtime = float(obj.last_modified)
                        info.update({"size_bytes": obj.size, "mtime": mtime})
                    files.append(info)
            return files
        except Exception as e:
            st.error(f"OSS 列表失败: {e}")
            return []

    def upload_file(self, filename: str, content: bytes) -> bool:
        key = f"{self.kb_prefix}/{filename}"
        try:
            # 检查是否存在
            self.oss_bucket.head_object(key)
            return False
        except NoSuchKey:
            pass
        try:
            self.oss_bucket.put_object(key, content)
            self.index_loaded = False
            return True
        except Exception:
            return False

    def delete_file(self, filename: str):
        key = f"{self.kb_prefix}/{filename}"
        try:
            self.oss_bucket.delete_object(key)
        except Exception:
            pass
        self.vector_store.remove_file(filename)
        remaining_files = self.get_file_list(with_metadata=True)
        self.vector_store.update_metadata(remaining_files)
        self.index_loaded = len(remaining_files) > 0 and self.vector_store.load_index()

    # ---------- 文档解析（从 OSS 读取字节流） ----------
    def _read_file_from_oss(self, filename: str) -> bytes:
        key = f"{self.kb_prefix}/{filename}"
        return self.oss_bucket.get_object(key).read()

    def _extract_text_from_bytes(self, filename: str, content_bytes: bytes) -> str:
        ext = filename.split('.')[-1].lower()
        # 对于需要文件路径的解析器，写入临时文件
        if ext in ['docx', 'doc', 'pdf', 'xlsx', 'xls', 'pptx', 'xmind']:
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                tmp.write(content_bytes)
                tmp_path = tmp.name
            try:
                return self._extract_text_from_path(tmp_path, ext)
            finally:
                os.unlink(tmp_path)
        else:
            # txt, csv, json, md 直接解码
            if ext in ['txt', 'md', 'json']:
                return content_bytes.decode('utf-8', errors='ignore')
            elif ext == 'csv':
                df = pd.read_csv(io.BytesIO(content_bytes))
                return df.to_string()
            else:
                return f"暂无法识别该类型 (.{ext})，请联系管理员"

    def _extract_text_from_path(self, filepath: str, ext: str) -> str:
        parsers = {
            "docx": lambda p: '\n'.join([para.text for para in docx.Document(p).paragraphs]) if HAS_DOCX else "请安装 python-docx",
            "doc": lambda p: '\n'.join([para.text for para in docx.Document(p).paragraphs]) if HAS_DOCX else "请安装 python-docx",
            "pdf": lambda p: '\n'.join([page.extract_text() for page in PyPDF2.PdfReader(p).pages]) if HAS_PDF else "请安装 PyPDF2",
            "xlsx": lambda p: self._parse_excel(p),
            "xls": lambda p: self._parse_excel(p),
            "pptx": lambda p: '\n'.join([shape.text for slide in Presentation(p).slides for shape in slide.shapes if hasattr(shape, "text")]) if HAS_PPTX else "请安装 python-pptx",
            "xmind": lambda p: self._parse_xmind(p),
        }
        if ext in parsers:
            try:
                return parsers[ext](filepath)
            except Exception as e:
                return f"提取失败: {e}"
        return f"暂无法识别该类型 (.{ext})，请联系管理员"

    @staticmethod
    def _parse_excel(filepath: str) -> str:
        if not HAS_EXCEL:
            return "请安装 openpyxl"
        wb = load_workbook(filepath, data_only=True)
        all_text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join([str(cell) for cell in row if cell is not None])
                if row_text:
                    all_text.append(row_text)
        return '\n'.join(all_text)

    @staticmethod
    def _parse_xmind(filepath: str) -> str:
        if not HAS_XMIND:
            return "请安装 xmindparser"
        data = xmindparser.xmind_to_dict(filepath)
        def extract(node):
            texts = []
            if isinstance(node, dict):
                if 'title' in node:
                    texts.append(node['title'])
                if 'note' in node and node['note']:
                    texts.append(node['note'])
                if 'topics' in node:
                    for sub in node['topics']:
                        texts.extend(extract(sub))
            elif isinstance(node, list):
                for item in node:
                    texts.extend(extract(item))
            return texts
        return '\n'.join(extract(data))

    # ---------- 分块与构建 ----------
    def _chunk_text(self, text: str, source: str) -> List[Dict]:
        prefixed_text = f"[文件：{source}]\n{text}"
        chunks = []
        paragraphs = prefixed_text.split('\n')
        current = ""
        chunk_id = 0
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            if len(para) > self.chunk_size:
                sentences = re.split(r'[。！？；]', para)
                for sentence in sentences:
                    sentence = sentence.strip()
                    if not sentence:
                        continue
                    if len(current) + len(sentence) <= self.chunk_size:
                        current += sentence + "。"
                    else:
                        if current:
                            chunks.append({"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
                            chunk_id += 1
                        overlap = current[-self.chunk_overlap:] if current else ""
                        current = overlap + sentence + "。"
            else:
                if len(current) + len(para) <= self.chunk_size:
                    current += para + "\n"
                else:
                    if current:
                        chunks.append({"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
                        chunk_id += 1
                    overlap = current[-self.chunk_overlap:] if current else ""
                    current = overlap + para + "\n"
        if current:
            chunks.append({"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
        return chunks

    def build_knowledge_base(self, progress_callback: Optional[Callable] = None, force=False) -> Dict:
        try:
            files = self.get_file_list(with_metadata=True)
            if not files:
                return {"status": "error", "message": "无文件"}

            if self.index_loaded:
                built_files = self.vector_store.get_built_files()
            else:
                if self.vector_store.load_index():
                    self.index_loaded = True
                    built_files = self.vector_store.get_built_files()
                else:
                    built_files = set()

            files_to_build = [f for f in files if f["name"] not in built_files]

            if not force and not files_to_build:
                stats = self.vector_store.get_stats()
                return {"status": "success", "message": f"所有文档已构建（共 {stats['documents']} 个知识片段）", "chunks": stats['documents']}

            # 增量构建
            if built_files and not force:
                new_chunks = []
                total_to_build = len(files_to_build)
                for i, f in enumerate(files_to_build):
                    if progress_callback:
                        progress_callback(i/total_to_build, f"处理 {f['name']} (解析中...)")
                    content_bytes = self._read_file_from_oss(f["name"])
                    text = self._extract_text_from_bytes(f["name"], content_bytes)
                    if text.startswith("请安装") or text.startswith("暂无法识别") or text.startswith("提取失败"):
                        continue
                    chunks = self._chunk_text(text, f["name"])
                    new_chunks.extend(chunks)
                    if progress_callback:
                        progress_callback(i/total_to_build, f"处理 {f['name']} → 生成 {len(chunks)} 个片段")

                if not new_chunks:
                    return {"status": "error", "message": "没有可构建的有效内容，请检查文件格式"}

                all_docs = self.vector_store.documents + [
                    {"id": i+len(self.vector_store.documents), "text": c["content"], "metadata": c["metadata"]}
                    for i, c in enumerate(new_chunks)
                ]
                texts = [doc["text"] for doc in all_docs]
                self.vector_store._build_vocab_and_idf(texts)
                vectors = [self.vector_store._text_to_vector(t) for t in texts]
                self.vector_store.vectors = np.array(vectors)
                self.vector_store.documents = all_docs
                # 保存到本地并上传 OSS
                with open(self.vector_store.index_file, 'w', encoding='utf-8') as f:
                    json.dump({
                        "documents": self.vector_store.documents,
                        "vocab": self.vector_store.vocab,
                        "idf": self.vector_store.idf
                    }, f, ensure_ascii=False)
                np.save(self.vector_store.vectors_file, self.vector_store.vectors)
                self.vector_store._upload_to_oss(self.vector_store.index_file, self.vector_store._oss_key("index.json"))
                self.vector_store._upload_to_oss(self.vector_store.vectors_file, self.vector_store._oss_key("vectors.npy"))
                files_info = [{"name": f["name"], "size_bytes": f["size_bytes"], "mtime": f["mtime"]} for f in files]
                self.vector_store.update_metadata(files_info)
                self.index_loaded = True
                return {"status": "success", "message": f"增量构建成功，新增 {len(files_to_build)} 个文件，{len(new_chunks)} 个知识片段", "files": len(files), "chunks": len(all_docs)}
            else:
                # 全量构建
                all_chunks = []
                missing_deps = set()
                unsupported_files = []
                total_files = len(files)
                for i, f in enumerate(files):
                    if progress_callback:
                        progress_callback(i/total_files, f"处理 {f['name']} (解析中...)")
                    content_bytes = self._read_file_from_oss(f["name"])
                    text = self._extract_text_from_bytes(f["name"], content_bytes)
                    if text.startswith("请安装"):
                        missing_deps.add(text)
                        continue
                    if text.startswith("暂无法识别该类型"):
                        unsupported_files.append(f"{f['name']}: {text}")
                        continue
                    if text.startswith("提取失败"):
                        unsupported_files.append(f"{f['name']}: {text}")
                        continue
                    chunks = self._chunk_text(text, f["name"])
                    all_chunks.extend(chunks)
                    if progress_callback:
                        progress_callback(i/total_files, f"处理 {f['name']} → 生成 {len(chunks)} 个片段")
                if missing_deps:
                    st.warning("\n".join(missing_deps))
                if unsupported_files:
                    for err in unsupported_files:
                        st.error(err)
                    if not all_chunks:
                        return {"status": "error", "message": "所有文件均无法解析，请检查依赖或文件格式"}
                if progress_callback:
                    progress_callback(0.95, "构建向量索引...")
                files_info = [{"name": f["name"], "size_bytes": f["size_bytes"], "mtime": f["mtime"]} for f in files]
                ok = self.vector_store.build_index(all_chunks, files_info)
                if ok:
                    self.index_loaded = True
                    return {"status": "success", "message": f"构建成功，共 {len(files)} 个文件，{len(all_chunks)} 个知识片段", "files": len(files), "chunks": len(all_chunks)}
                else:
                    return {"status": "error", "message": "索引失败"}
        except Exception as e:
            st.error(f"构建知识库时发生未预期的错误: {str(e)}")
            return {"status": "error", "message": str(e)}

    def search_knowledge(self, query: str, top_k=5):
        if not self.index_loaded and not self.vector_store.load_index():
            return []
        return self.vector_store.search(query, top_k)

    def get_knowledge_context(self, query: str, max_chunks=5) -> str:
        results = self.search_knowledge(query, max_chunks)
        if not results:
            return ""
        parts = []
        for i, r in enumerate(results, 1):
            parts.append(f"【知识片段 {i}】来源：{r['metadata'].get('source','未知')} 相关度：{r['similarity']:.2%}\n{r['content']}")
        return "\n\n---\n\n".join(parts)

    def get_built_files_safe(self) -> set:
        if self.index_loaded:
            return self.vector_store.get_built_files()
        else:
            if self.vector_store.load_index():
                self.index_loaded = True
                return self.vector_store.get_built_files()
            return set()

    # 然后原样复制原代码中 EnhancedKnowledgeBase 的所有方法，包括：
    # _auto_load, refresh_index, get_file_list, upload_file, delete_file,
    # _read_file_from_oss, _extract_text_from_bytes, _extract_text_from_path,
    # _parse_excel, _parse_xmind, _chunk_text, build_knowledge_base,
    # search_knowledge, get_knowledge_context, get_built_files_safe