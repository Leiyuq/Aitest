"""
测试用例智构系统 - 轻量级 RAG（NumPy TF‑IDF）
启动方式（允许局域网访问）：
    streamlit run pros.py --server.address 0.0.0.0
"""
import streamlit as st
st.set_page_config(page_title="测试用例智构系统", layout="wide")

import os
import re
import json
import pickle
import io
import time
import warnings
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Callable
from dataclasses import dataclass
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError

import numpy as np
import pandas as pd
from openai import OpenAI

# 可选依赖：用于解析特殊格式
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from openpyxl import load_workbook
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import xmindparser
    HAS_XMIND = True
except ImportError:
    HAS_XMIND = False

try:
    import jieba
    USE_JIEBA = True
except ImportError:
    USE_JIEBA = False
    jieba = None

warnings.filterwarnings('ignore')


# ====================== 配置 ======================
@dataclass
class ModelConfig:
    name: str
    api_key: str
    base_url: str
    model: str

class AppConfig:
    PAGE_TITLE = "测试用例智构系统"
    LAYOUT = "wide"
    BASE_ROOT = "user_data"
    ALLOWED_FILE_TYPES = {
        "txt", "csv", "docx", "doc", "pdf", "xlsx", "xls", "pptx", "xmind", "md"
    }
    RAG_TOP_K = 5
    RAG_SIMILARITY_THRESHOLD = 0.01
    API_TIMEOUT = 60
    API_MAX_RETRIES = 2
    API_RETRY_DELAY = 2

    MODELS = {
        "qwen": ModelConfig(
            name="通义千问",
            api_key="sk-66b40867237b4e589c81c0255ff94b36",
            base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
            model="qwen-plus"
        ),
        "local": ModelConfig(
            name="本地离线",
            api_key="none",
            base_url="none",
            model="local"
        )
    }

    @classmethod
    def get_model_list(cls):
        return [f"{m.name}({k})" for k, m in cls.MODELS.items() if m.api_key or k == "local"]


# ====================== 轻量级 TF‑IDF 向量检索 ======================
class SimpleVectorStore:
    def __init__(self, store_dir: str):
        self.store_dir = store_dir
        self.index_file = os.path.join(store_dir, "index.json")
        self.vectors_file = os.path.join(store_dir, "vectors.npy")
        self.metadata_file = os.path.join(store_dir, "metadata.pkl")
        self.documents = []
        self.vocab = {}
        self.idf = {}
        self.vectors = None
        Path(store_dir).mkdir(parents=True, exist_ok=True)

    @staticmethod
    def _tokenize(text: str) -> List[str]:
        text = text.lower()
        if USE_JIEBA and jieba:
            words = jieba.lcut(text)
            return [w for w in words if len(w.strip()) > 1 or w.isalnum()]
        else:
            return re.findall(r'[\u4e00-\u9fa5]{2,}|[a-z]{2,}', text)

    @staticmethod
    def _clean_query(query: str) -> str:
        """移除编号、JIRA单号、日期等噪声"""
        query = re.sub(r'\b[A-Za-z0-9]+[_-][A-Za-z0-9]+\b', '', query)
        query = re.sub(r'\b[A-Z]{2,}[0-9]+\b', '', query)
        query = re.sub(r'\b[A-Z]+-\d+\b', '', query)
        query = re.sub(r'\b\d{4}[-/]\d{1,2}[-/]\d{1,2}\b', '', query)
        query = re.sub(r'\b\d{1,2}[-/]\d{1,2}[-/]\d{4}\b', '', query)
        query = re.sub(r'\b\d{8}\b', '', query)
        query = re.sub(r'\s+', ' ', query).strip()
        return query

    def _build_vocab_and_idf(self, all_texts: List[str]):
        doc_count = {}
        for text in all_texts:
            unique_words = set(self._tokenize(text))
            for w in unique_words:
                doc_count[w] = doc_count.get(w, 0) + 1

        total_docs = len(all_texts)
        self.vocab = {}
        self.idf = {}
        for w, df in doc_count.items():
            if 0.01 < df / total_docs < 0.8:
                idx = len(self.vocab)
                self.vocab[w] = idx
                self.idf[w] = np.log((total_docs + 1) / (df + 0.5))

    def _text_to_vector(self, text: str) -> np.ndarray:
        words = self._tokenize(text)
        if not words or not self.vocab:
            return np.zeros(len(self.vocab))
        tf = {}
        for w in words:
            if w in self.vocab:
                tf[w] = tf.get(w, 0) + 1
        vec = np.zeros(len(self.vocab))
        for w, cnt in tf.items():
            idx = self.vocab[w]
            tf_val = np.log1p(cnt / len(words))
            vec[idx] = tf_val * self.idf.get(w, 1.0)
        norm = np.linalg.norm(vec)
        if norm > 0:
            vec /= norm
        return vec

    def build_index(self, documents: List[Dict], files_info: List[Dict] = None) -> bool:
        try:
            if not documents:
                return False
            self.documents = [{"id": i, "text": d["content"], "metadata": d.get("metadata", {})}
                              for i, d in enumerate(documents)]
            texts = [d["text"] for d in self.documents]
            self._build_vocab_and_idf(texts)
            vectors = [self._text_to_vector(t) for t in texts]
            self.vectors = np.array(vectors)

            with open(self.index_file, 'w', encoding='utf-8') as f:
                json.dump({
                    "documents": self.documents,
                    "vocab": self.vocab,
                    "idf": self.idf
                }, f, ensure_ascii=False)
            np.save(self.vectors_file, self.vectors)
            if files_info:
                with open(self.metadata_file, 'wb') as f:
                    pickle.dump({"files_info": files_info, "updated": datetime.now().isoformat()}, f)
            return True
        except Exception as e:
            st.error(f"构建索引失败: {e}")
            return False

    def load_index(self) -> bool:
        try:
            if os.path.exists(self.index_file) and os.path.exists(self.vectors_file):
                with open(self.index_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    self.documents = data["documents"]
                    self.vocab = data["vocab"]
                    self.idf = data["idf"]
                self.vectors = np.load(self.vectors_file)
                return True
        except Exception as e:
            st.warning(f"加载索引失败: {e}")
        return False

    def search(self, query: str, top_k: int = 5) -> List[Dict]:
        if self.vectors is None or len(self.documents) == 0:
            return []
        cleaned_query = self._clean_query(query)
        if not cleaned_query:
            cleaned_query = query
        q_vec = self._text_to_vector(cleaned_query)
        scores = np.dot(self.vectors, q_vec)
        top_idx = np.argsort(scores)[-top_k:][::-1]
        results = []
        for idx in top_idx:
            sim = float(scores[idx])
            if sim >= AppConfig.RAG_SIMILARITY_THRESHOLD:
                doc = self.documents[idx]
                results.append({
                    "content": doc["text"],
                    "metadata": doc["metadata"],
                    "similarity": sim
                })
        return results

    def is_index_valid(self, current_files_info: List[Dict]) -> bool:
        if not os.path.exists(self.metadata_file):
            return False
        with open(self.metadata_file, 'rb') as f:
            meta = pickle.load(f)
        saved = {f["name"]: f for f in meta["files_info"]}
        curr = {f["name"]: f for f in current_files_info}
        if set(saved.keys()) != set(curr.keys()):
            return False
        for name, cf in curr.items():
            sf = saved[name]
            if sf.get("size_bytes") != cf.get("size_bytes") or sf.get("mtime") != cf.get("mtime"):
                return False
        return True

    def get_built_files(self) -> set:
        if not self.documents:
            return set()
        return {doc["metadata"].get("source") for doc in self.documents if "source" in doc["metadata"]}

    def get_stats(self) -> Dict:
        return {
            "documents": len(self.documents),
            "vocab_size": len(self.vocab),
            "has_index": self.vectors is not None
        }


# ====================== 知识库封装 ======================
class EnhancedKnowledgeBase:
    def __init__(self, user_dir: str):
        self.user_dir = user_dir
        self.kb_dir = os.path.join(user_dir, "knowledge_base")
        self.vector_store = SimpleVectorStore(os.path.join(user_dir, "vector_store"))
        self.chunk_size = 800
        self.chunk_overlap = 200
        Path(self.kb_dir).mkdir(parents=True, exist_ok=True)
        self.index_loaded = False
        self._auto_load()

    def _auto_load(self):
        files = self.get_file_list(with_metadata=True)
        if files and self.vector_store.is_index_valid(files):
            if self.vector_store.load_index():
                self.index_loaded = True

    def refresh_index(self) -> Dict:
        self.index_loaded = False
        files = self.get_file_list(with_metadata=True)
        if not files:
            return {"status": "error", "message": "知识库无文件"}
        if self.vector_store.is_index_valid(files):
            if self.vector_store.load_index():
                self.index_loaded = True
                stats = self.vector_store.get_stats()
                return {
                    "status": "success",
                    "message": f"已加载索引，包含 {stats['documents']} 个文档片段，词表大小 {stats['vocab_size']}"
                }
        return {"status": "error", "message": "未找到有效索引，请先构建知识库"}

    def get_file_list(self, with_metadata=False):
        files = []
        if os.path.exists(self.kb_dir):
            for name in os.listdir(self.kb_dir):
                ext = name.split('.')[-1].lower()
                if ext in AppConfig.ALLOWED_FILE_TYPES:
                    path = os.path.join(self.kb_dir, name)
                    stat = os.stat(path)
                    info = {"name": name, "size": f"{stat.st_size/1024:.1f}KB"}
                    if with_metadata:
                        info.update({"size_bytes": stat.st_size, "mtime": stat.st_mtime})
                    files.append(info)
        return files

    def upload_file(self, filename: str, content: bytes):
        path = os.path.join(self.kb_dir, filename)
        with open(path, "wb") as f:
            f.write(content)
        self.index_loaded = False

    def delete_file(self, filename: str):
        path = os.path.join(self.kb_dir, filename)
        if os.path.exists(path):
            os.remove(path)
        self.index_loaded = False

    # ---------- 文档解析 ----------
    @staticmethod
    def _extract_text(filepath: str) -> str:
        ext = filepath.split('.')[-1].lower()
        parsers = {
            "txt": lambda p: open(p, 'r', encoding='utf-8').read(),
            "md": lambda p: open(p, 'r', encoding='utf-8').read(),
            "csv": lambda p: pd.read_csv(p).to_string(),
            "json": lambda p: json.dumps(json.load(open(p, 'r', encoding='utf-8')), ensure_ascii=False),
            "docx": lambda p: '\n'.join([para.text for para in docx.Document(p).paragraphs]) if HAS_DOCX else "请安装 python-docx",
            "doc": lambda p: '\n'.join([para.text for para in docx.Document(p).paragraphs]) if HAS_DOCX else "请安装 python-docx",
            "pdf": lambda p: '\n'.join([page.extract_text() for page in PyPDF2.PdfReader(p).pages]) if HAS_PDF else "请安装 PyPDF2",
            "xlsx": lambda p: EnhancedKnowledgeBase._parse_excel(p),
            "xls": lambda p: EnhancedKnowledgeBase._parse_excel(p),
            "pptx": lambda p: '\n'.join([shape.text for slide in Presentation(p).slides for shape in slide.shapes if hasattr(shape, "text")]) if HAS_PPTX else "请安装 python-pptx",
            "xmind": lambda p: EnhancedKnowledgeBase._parse_xmind(p),
        }
        if ext in parsers:
            try:
                return parsers[ext](filepath)
            except Exception as e:
                return f"提取失败: {e}"
        return f"暂无法识别该类型 (.{ext})，请联系管理员"

    @staticmethod
    def _parse_excel(filepath: str) -> str:
        if not HAS_EXCEL:
            return "请安装 openpyxl"
        wb = load_workbook(filepath, data_only=True)
        all_text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join([str(cell) for cell in row if cell is not None])
                if row_text:
                    all_text.append(row_text)
        return '\n'.join(all_text)

    @staticmethod
    def _parse_xmind(filepath: str) -> str:
        if not HAS_XMIND:
            return "请安装 xmindparser"
        data = xmindparser.xmind_to_dict(filepath)
        def extract(node):
            texts = []
            if isinstance(node, dict):
                if 'title' in node:
                    texts.append(node['title'])
                if 'topics' in node:
                    for sub in node['topics']:
                        texts.extend(extract(sub))
            elif isinstance(node, list):
                for item in node:
                    texts.extend(extract(item))
            return texts
        return '\n'.join(extract(data))

    # ---------- 分块与构建（文件名加入内容前缀） ----------
    def _chunk_text(self, text: str, source: str) -> List[Dict]:
        # 将文件名作为前缀加入，使文件名也能被检索到
        prefixed_text = f"[文件：{source}]\n{text}"
        chunks = []
        paragraphs = prefixed_text.split('\n')
        current = ""
        chunk_id = 0
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            if len(para) > self.chunk_size:
                sentences = re.split(r'[。！？；]', para)
                for sentence in sentences:
                    sentence = sentence.strip()
                    if not sentence:
                        continue
                    if len(current) + len(sentence) <= self.chunk_size:
                        current += sentence + "。"
                    else:
                        if current:
                            chunks.append({"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
                            chunk_id += 1
                        overlap = current[-self.chunk_overlap:] if current else ""
                        current = overlap + sentence + "。"
            else:
                if len(current) + len(para) <= self.chunk_size:
                    current += para + "\n"
                else:
                    if current:
                        chunks.append({"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
                        chunk_id += 1
                    overlap = current[-self.chunk_overlap:] if current else ""
                    current = overlap + para + "\n"
        if current:
            chunks.append({"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
        return chunks

    def build_knowledge_base(self, progress_callback: Optional[Callable] = None, force=False) -> Dict:
        try:
            files = self.get_file_list(with_metadata=True)
            if not files:
                return {"status": "error", "message": "无文件"}
            if not force and self.index_loaded:
                if self.vector_store.is_index_valid(files):
                    stats = self.vector_store.get_stats()
                    return {"status": "success", "message": f"当前文档已构建无需重复（已加载 {stats['documents']} 个知识片段）", "chunks": stats['documents'], "loaded_from_cache": True}
            all_chunks = []
            missing_deps = set()
            unsupported_files = []
            total_files = len(files)
            for i, f in enumerate(files):
                if progress_callback:
                    progress_callback(i/total_files, f"处理 {f['name']} (解析中...)")
                text = self._extract_text(os.path.join(self.kb_dir, f["name"]))
                if text.startswith("请安装"):
                    missing_deps.add(text)
                    continue
                if text.startswith("暂无法识别该类型"):
                    unsupported_files.append(f"{f['name']}: {text}")
                    continue
                if text.startswith("提取失败"):
                    unsupported_files.append(f"{f['name']}: {text}")
                    continue
                chunks = self._chunk_text(text, f["name"])
                all_chunks.extend(chunks)
                if progress_callback:
                    progress_callback(i/total_files, f"处理 {f['name']} → 生成 {len(chunks)} 个片段")
            if missing_deps:
                st.warning("\n".join(missing_deps))
            if unsupported_files:
                for err in unsupported_files:
                    st.error(err)
                if not all_chunks:
                    return {"status": "error", "message": "所有文件均无法解析，请检查依赖或文件格式"}
            if progress_callback:
                progress_callback(0.95, "构建向量索引...")
            files_info = [{"name": f["name"], "size_bytes": f["size_bytes"], "mtime": f["mtime"]} for f in files]
            ok = self.vector_store.build_index(all_chunks, files_info)
            if ok:
                self.index_loaded = True
                return {"status": "success", "message": f"构建成功，共 {len(files)} 个文件，{len(all_chunks)} 个知识片段", "files": len(files), "chunks": len(all_chunks)}
            else:
                return {"status": "error", "message": "索引失败"}
        except Exception as e:
            st.error(f"构建知识库时发生未预期的错误: {str(e)}")
            return {"status": "error", "message": str(e)}

    def search_knowledge(self, query: str, top_k=5):
        if not self.index_loaded and not self.vector_store.load_index():
            return []
        return self.vector_store.search(query, top_k)

    def get_knowledge_context(self, query: str, max_chunks=5) -> str:
        results = self.search_knowledge(query, max_chunks)
        if not results:
            return ""
        parts = []
        for i, r in enumerate(results, 1):
            parts.append(f"【知识片段 {i}】来源：{r['metadata'].get('source','未知')} 相关度：{r['similarity']:.2%}\n{r['content']}")
        return "\n\n---\n\n".join(parts)

    def get_built_files_safe(self) -> set:
        if self.index_loaded:
            return self.vector_store.get_built_files()
        else:
            if self.vector_store.load_index():
                self.index_loaded = True
                return self.vector_store.get_built_files()
            return set()


# ====================== 用户管理 ======================
class UserRepository:
    @staticmethod
    def _get_user_file_path() -> str:
        Path(AppConfig.BASE_ROOT).mkdir(parents=True, exist_ok=True)
        return os.path.join(AppConfig.BASE_ROOT, "users.csv")

    @staticmethod
    def load_all() -> pd.DataFrame:
        user_file = UserRepository._get_user_file_path()
        if not os.path.exists(user_file):
            df = pd.DataFrame({
                "username": ["admin", "user1"],
                "password": ["123456", "123456"],
                "role": ["admin", "user"],
                "status": ["正常", "正常"],
                "created_at": [datetime.now().strftime("%Y-%m-%d %H:%M")] * 2
            })
            df.to_csv(user_file, index=False, encoding="utf-8")
            return df
        try:
            df = pd.read_csv(user_file, encoding="utf-8")
            required_cols = {"username", "password", "role", "status", "created_at"}
            if not required_cols.issubset(df.columns):
                return UserRepository.load_all()
            for col in ["username", "password", "role"]:
                if col in df.columns:
                    df[col] = df[col].astype(str).str.strip()
            if "status" not in df.columns:
                df["status"] = "正常"
            df["status"] = df["status"].astype(str).str.strip()
            return df
        except Exception as e:
            st.error(f"读取用户文件失败: {e}")
            if os.path.exists(user_file):
                backup = user_file + ".bak"
                os.rename(user_file, backup)
                st.warning(f"已备份损坏的用户文件至 {backup}")
            return UserRepository.load_all()

    @staticmethod
    def save_all(df: pd.DataFrame) -> None:
        user_file = UserRepository._get_user_file_path()
        df.to_csv(user_file, index=False, encoding="utf-8")

    @staticmethod
    def authenticate(username: str, password: str) -> Optional[Dict[str, str]]:
        df = UserRepository.load_all()
        user = df[(df["username"] == username) & (df["password"] == password)]
        if not user.empty and user.iloc[0]["status"] == "正常":
            return {"username": user.iloc[0]["username"], "role": user.iloc[0]["role"]}
        return None

    @staticmethod
    def create(username: str, password: str) -> bool:
        df = UserRepository.load_all()
        if username in df["username"].values:
            return False
        new_row = {
            "username": username,
            "password": password,
            "role": "user",
            "status": "正常",
            "created_at": datetime.now().strftime("%Y-%m-%d %H:%M")
        }
        df = pd.concat([df, pd.DataFrame([new_row])], ignore_index=True)
        UserRepository.save_all(df)
        return True


# ====================== LLM 服务 ======================
class LLMService:
    def __init__(self, model_key: str):
        ident = model_key.split("(")[-1].rstrip(")")
        self.config = AppConfig.MODELS.get(ident, AppConfig.MODELS["local"])

    def generate_cases(self, prompt: str, context: str) -> Dict:
        if self.config.model != "local" and not self.config.api_key:
            return {"status": "error", "message": f"未配置 {self.config.name} API Key"}
        if self.config.model == "local":
            return {"status": "success", "content": self._local_generate(), "message": "本地生成"}

        for attempt in range(AppConfig.API_MAX_RETRIES):
            try:
                client = OpenAI(api_key=self.config.api_key, base_url=self.config.base_url, timeout=AppConfig.API_TIMEOUT)
                messages = [
                    {"role": "system", "content": self._system_prompt()},
                    {"role": "user", "content": f"需求：{prompt}\n\n知识库：\n{context}"}
                ]
                with ThreadPoolExecutor(max_workers=1) as executor:
                    future = executor.submit(
                        lambda: client.chat.completions.create(
                            model=self.config.model, messages=messages, temperature=0.3, max_tokens=2000
                        )
                    )
                    resp = future.result(timeout=AppConfig.API_TIMEOUT)
                    return {"status": "success", "content": resp.choices[0].message.content.strip(), "message": "生成成功"}
            except FuturesTimeoutError:
                if attempt < AppConfig.API_MAX_RETRIES - 1:
                    time.sleep(AppConfig.API_RETRY_DELAY)
                    continue
                return {"status": "error", "message": "请求超时，请稍后重试"}
            except Exception as e:
                if attempt == AppConfig.API_MAX_RETRIES - 1:
                    return {"status": "error", "message": f"生成失败: {str(e)}"}
                time.sleep(AppConfig.API_RETRY_DELAY)
        return {"status": "error", "message": "未知错误"}

    @staticmethod
    def _system_prompt():
        return """你是一名资深测试工程师。请根据给定的需求和知识库，生成规范、可执行的测试用例。输出需严格遵循以下格式，不输出任何额外解释或标记。

    ## 输出格式
    用例ID： TC001
    JIRA单号： DEMO-001
    用例名称： [动词+功能点，如：用户登录成功]
    前置条件： [明确的环境、数据或状态；多条用分号分隔]
    测试步骤： 1. 操作A; 2. 操作B; 3. 操作C
    预期结果： 1. 结果A; 2. 结果B; 3. 结果C

    多个用例之间用空行分隔。每条用例的“用例ID”和“JIRA单号”字段仅输出值，不加方括号或引号。

    ## 用例设计要求
    - 覆盖正向流程、异常场景、边界值、权限/状态依赖等。
    - 测试步骤应包含具体操作数据（如输入值、点击元素）。
    - 预期结果应可客观验证，避免模糊描述（如“系统正常”）。
    - 若需求隐含约束（如未登录、网络异常），请设计对应负向用例。

    ## 示例（仅供参考，不要输出到结果中）
    用例ID： TC002
    JIRA单号： DEMO-002
    用例名称： 登录密码错误
    前置条件： 已注册用户test01，密码正确为123456；登录页已打开
    测试步骤： 1. 输入用户名test01; 2. 输入错误密码654321; 3. 点击登录按钮
    预期结果： 1-2. 输入框正常接收；3. 页面提示“用户名或密码错误”，未跳转"""

    @staticmethod
    def _local_generate():
        return """用例ID： TC001
JIRA单号： DEMO-001
用例名称： 功能验证
前置条件： 环境正常
测试步骤： 1.操作;2.验证
预期结果： 1.成功;2.正常"""


# ====================== 用例解析与导出 ======================
class TestCaseService:
    @staticmethod
    def parse(content: str, jira_codes: List[str]) -> List[Dict]:
        cases = []
        blocks = re.split(r'\n\s*\n', content.strip())
        jira_idx = 0
        for block in blocks:
            if not any(k in block for k in ["用例ID", "用例名称", "测试步骤"]):
                continue
            case = {"JIRA单号": "", "用例ID": "", "用例名称": "", "前置条件": "", "测试步骤": "", "预期结果": ""}
            patterns = {
                "用例ID": r"用例ID[:：]\s*([^\n]+)",
                "JIRA单号": r"JIRA单号[:：]\s*([^\n]+)",
                "用例名称": r"用例名称[:：]\s*([^\n]+)",
                "前置条件": r"前置条件[:：]\s*([^\n]+(?:\n(?!用例|JIRA|测试|预期)[^\n]+)*)",
                "测试步骤": r"测试步骤[:：]\s*([^\n]+(?:\n(?!用例|JIRA|前置|预期)[^\n]+)*)",
                "预期结果": r"预期结果[:：]\s*([^\n]+(?:\n(?!用例|JIRA|前置|测试)[^\n]+)*)"
            }
            for field, pat in patterns.items():
                m = re.search(pat, block, re.DOTALL)
                if m:
                    value = m.group(1).strip()
                    if field in ["用例ID", "JIRA单号"]:
                        value = re.sub(r'^\[|\]$', '', value)
                    case[field] = re.sub(r'\n+', ' ', value)
            if not case["JIRA单号"] and jira_codes:
                case["JIRA单号"] = jira_codes[jira_idx % len(jira_codes)]
                jira_idx += 1
            if case["用例ID"]:
                cases.append(case)
        return cases

class ExportService:
    @staticmethod
    def to_csv(cases: List[Dict]) -> bytes:
        df = pd.DataFrame(cases)
        cols = ["JIRA单号", "用例ID", "用例名称", "前置条件", "测试步骤", "预期结果"]
        df = df[[c for c in cols if c in df.columns]]
        buf = io.BytesIO()
        df.to_csv(buf, index=False, encoding='utf-8-sig')
        return buf.getvalue()


# ====================== 视图层 ======================
class LoginView:
    @staticmethod
    def render() -> bool:
        st.markdown("<h1 style='text-align: center;'>测试用例智构系统</h1>", unsafe_allow_html=True)
        col1, col2, col3 = st.columns([0.3, 0.4, 0.3])
        with col2:
            tab1, tab2 = st.tabs(["登录", "注册"])
            with tab1:
                with st.form("login_form"):
                    st.text_input("用户名", key="login_username")
                    st.text_input("密码", type="password", key="login_password")
                    submitted = st.form_submit_button("登录")
                    if submitted:
                        username = st.session_state.get("login_username", "")
                        password = st.session_state.get("login_password", "")
                        user = UserRepository.authenticate(username, password)
                        if user:
                            st.session_state.user = user
                            st.rerun()
                        else:
                            st.error("用户名或密码错误")
            with tab2:
                with st.form("register_form"):
                    st.text_input("用户名", key="reg_username")
                    st.text_input("密码", type="password", key="reg_password")
                    submitted_reg = st.form_submit_button("注册")
                    if submitted_reg:
                        new_user = st.session_state.get("reg_username", "")
                        new_pwd = st.session_state.get("reg_password", "")
                        if UserRepository.create(new_user, new_pwd):
                            st.success("注册成功，请登录")
                        else:
                            st.error("用户名已存在")
        return "user" in st.session_state

class MainView:
    def __init__(self):
        self.user = st.session_state.user
        data_dir = os.path.join(AppConfig.BASE_ROOT, self.user["username"])
        Path(data_dir).mkdir(exist_ok=True)
        self.kb = EnhancedKnowledgeBase(data_dir)

    def render(self):
        st.title(f"欢迎，{self.user['username']}")
        with st.sidebar:
            if st.button("退出登录"):
                del st.session_state.user
                st.rerun()
            st.divider()
            model_key = st.selectbox("模型选择", AppConfig.get_model_list())
            st.subheader("文件管理")
            uploaded = st.file_uploader("上传文档", type=list(AppConfig.ALLOWED_FILE_TYPES))
            if uploaded and st.button("上传"):
                self.kb.upload_file(uploaded.name, uploaded.getvalue())
                st.success("上传成功")
                st.rerun()

        self._kb_panel()
        st.divider()
        self._gen_panel(model_key)

    def _kb_panel(self):
        st.subheader("知识库")
        files = self.kb.get_file_list()
        built_files = self.kb.get_built_files_safe()
        if files:
            for f in files:
                # 调整列宽：状态列0.5，文件名列4.5，删除按钮列0.8（左移）
                col1, col2, col3 = st.columns([0.5, 3, 0.8])
                is_built = f['name'] in built_files
                if is_built:
                    col1.write("✅已构建")
                else:
                    col1.write("⏳未构建")
                col2.write(f"📄 {f['name']} ({f['size']})")
                delete_key = f"del_{f['name']}"
                if is_built:
                    # 已构建文档：popover 二次确认
                    with col3.popover("删除", use_container_width=False):
                        st.write(f"确定要删除文档 **{f['name']}** 吗？")
                        st.write("删除后，该文档构建的索引也会被清除。")

                        # 调整为更紧凑的列宽比例，让取消按钮左移，靠近确认删除
                        col_confirm, col_cancel, _ = st.columns([1, 1, 2])

                        with col_confirm:
                            if st.button("确认删除", key=f"confirm_{delete_key}", type="primary"):
                                self._delete_file_and_rebuild(f['name'])
                                st.session_state.popover_file = None  # 清空状态
                                st.rerun()

                        with col_cancel:
                         if st.button("取消",key=f"cancel_{delete_key}"):
                            st.session_state.popover_file = None  # 清空状态
                            st.rerun()  # 轻微刷新，完全不影响体验
                else:
                    # 未构建文档：直接删除
                    if col3.button("删除", key=delete_key):
                        self._delete_file_and_rebuild(f['name'])
                        st.rerun()
        else:
            st.info("暂无文档，请先上传文件")

        col_btn1, col_btn2, col_help, _ = st.columns([1, 1, 0.3, 6])
        with col_btn1:
            if st.button("构建知识库", type="primary"):
                if not files:
                    st.warning("请先上传文件")
                else:
                    with st.status("正在构建知识库...", expanded=True) as status:
                        prog = st.progress(0)
                        stat = st.empty()
                        def cb(p, msg):
                            prog.progress(p)
                            stat.text(msg)
                        res = self.kb.build_knowledge_base(cb, force=False)
                        if res["status"] == "success":
                            status.update(state="complete")  # 清空label，避免显示“构建完成”
                            st.toast(res['message'], icon="✅")
                        else:
                            status.update(label="", state="error")
                            st.toast(res['message'], icon="❌")
                    time.sleep(2)
                    st.rerun()
        with col_btn2:
            if st.button("刷新索引", help="重新加载已有索引"):
                res = self.kb.refresh_index()
                if res["status"] == "success":
                    st.toast(res["message"], icon="🔄")
                else:
                    st.toast(res["message"], icon="⚠️")
                time.sleep(1)
                st.rerun()

    def _delete_file_and_rebuild(self, filename: str):
        """删除文件并重建索引"""
        self.kb.delete_file(filename)
        remaining_files = self.kb.get_file_list()
        if remaining_files:
            with st.spinner("正在更新知识库索引..."):
                res = self.kb.build_knowledge_base(force=True)
                if res["status"] == "success":
                    st.toast("索引已更新", icon="✅")
                else:
                    st.toast(f"索引更新失败: {res['message']}", icon="❌")
        else:
            self.kb.index_loaded = False

    def _gen_panel(self, model_key):
        st.subheader("生成测试用例")
        input_type = st.radio("输入方式", ["文本输入", "JIRA单号"], horizontal=True)
        if input_type == "文本输入":
            prompt = st.text_area("需求描述", height=150)
        else:
            jira = st.text_input("JIRA单号")
            prompt = st.text_area("需求描述", value=jira, height=150) if jira else ""

        col_check, col_help, _ = st.columns([1, 0.3, 6])
        with col_check:
            use_rag = st.checkbox("启用RAG知识库 ", value=True,
                                  help="从知识库中检索相关内容，提升生成质量")

        col_gen, col_clear, _ = st.columns([1, 1, 6])
        with col_gen:
            generate_clicked = st.button("生成测试用例", type="primary")
        with col_clear:
            clear_clicked = st.button("清空结果")

        if clear_clicked:
            if "cases" in st.session_state:
                del st.session_state.cases
            if "result" in st.session_state:
                del st.session_state.result
            st.rerun()

        if generate_clicked:
            if not prompt:
                st.warning("请输入需求描述")
                return
            with st.spinner("生成中..."):
                context = ""
                if use_rag and self.kb.get_file_list():
                    if not self.kb.index_loaded:
                        self.kb.refresh_index()
                    if self.kb.index_loaded:
                        with st.expander("RAG检索详情"):
                            results = self.kb.search_knowledge(prompt, top_k=5)
                            if results:
                                for r in results:
                                    st.write(f"相似度 {r['similarity']:.3f} | 来源 {r['metadata'].get('source','')}")
                                    st.caption(r['content'][:200])
                            else:
                                st.info("未检索到相关知识")
                        context = self.kb.get_knowledge_context(prompt, max_chunks=5)
                    else:
                        st.info("知识库尚未构建，将直接使用模型生成。如需检索知识，请先构建知识库。")
                llm = LLMService(model_key)
                resp = llm.generate_cases(prompt, context)
                if resp["status"] == "error":
                    st.error(resp["message"])
                else:
                    st.success(resp["message"])
                    jira_codes = re.findall(r'[A-Za-z0-9_]+-\d+', prompt)
                    cases = TestCaseService.parse(resp["content"], jira_codes)
                    if cases:
                        st.session_state.cases = cases
                        st.success(f"解析出 {len(cases)} 个测试用例")
                    else:
                        st.warning("解析失败，显示原始内容")
                        st.code(resp["content"])
        if "cases" in st.session_state:
            self._show_results()

    @staticmethod
    def _show_results():
        cases = st.session_state.cases
        if not cases:
            return
        st.subheader("测试用例")
        df = pd.DataFrame(cases)
        cols = ["JIRA单号", "用例ID", "用例名称", "前置条件", "测试步骤", "预期结果"]
        df = df[[c for c in cols if c in df.columns]]
        st.dataframe(df, use_container_width=True)
        csv = ExportService.to_csv(cases)
        st.download_button("导出 CSV", data=csv, file_name=f"cases_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv", mime="text/csv")


class AdminView:
    @staticmethod
    def render():
        st.title("用户管理")
        df = UserRepository.load_all()
        st.dataframe(df[["username", "role", "status"]])
        with st.expander("添加用户"):
            with st.form("add_user_form"):
                new_user = st.text_input("用户名")
                new_pwd = st.text_input("密码", type="password")
                submitted = st.form_submit_button("添加")
                if submitted:
                    if UserRepository.create(new_user, new_pwd):
                        st.success("添加成功")
                        st.rerun()
                    else:
                        st.error("用户已存在")


# ====================== 入口 ======================
def main():
    Path(AppConfig.BASE_ROOT).mkdir(exist_ok=True)
    if "user" not in st.session_state:
        if LoginView.render():
            st.rerun()
    else:
        role = st.session_state.user["role"]
        if role == "admin":
            page = st.sidebar.radio("导航", ["主页", "管理"])
            if page == "主页":
                MainView().render()
            else:
                AdminView.render()
        else:
            MainView().render()

if __name__ == "__main__":
    main()